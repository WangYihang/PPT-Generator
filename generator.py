#!/usr/bin/env python
# encoding: utf-8

import sys
from anytree.importer import IndentedStringImporter
from pptx import Presentation
from anytree import RenderTree
from anytree import LevelOrderIter

def ascii_tree(tree):
    for pre, fill, node in RenderTree(tree):
        print("%s%s" % (pre, node.name))

def create_slide(ppt, nodes):
    bullet_slide_layout = ppt.slide_layouts[1]
    slide = ppt.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = nodes[0].name
    body_shape = slide.shapes.placeholders
    body_shape = slide.shapes.placeholders
    for i in nodes[1:]:
        new_paragraph = body_shape[1].text_frame.add_paragraph()
        new_paragraph.text = i.name
        new_paragraph.level = 1

def walk(tree):
    slide = []
    for i in LevelOrderIter(tree, maxlevel=2):
        slide.append(i)
    return slide

def run(ppt, tree):
    if len(tree.children) == 0:
        return
    titles = walk(tree)
    for i in titles[1:]:
        t = walk(i)
        create_slide(ppt, t)
        run(ppt, i)

def main():
    if len(sys.argv) != 3:
        print("Usage: ")
        print("\tpython %s [INPUT] [OUTPUT]" % (sys.argv[0]))
        print("Example: ")
        print("\tpython %s input.txt output.pptx" % (sys.argv[0]))
        exit(1)
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    ppt = Presentation()
    data = open(input_file).read()
    importer = IndentedStringImporter()
    root = importer.import_(data)
    ascii_tree(root)
    run(ppt, root)
    ppt.save(output_file)

if __name__ == '__main__':
    main()

